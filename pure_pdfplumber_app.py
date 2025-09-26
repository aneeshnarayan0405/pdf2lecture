# pure_pdfplumber_app.py
import os
import pdfplumber
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont
import textwrap

def simple_pdf_to_lecture(pdf_path, output_dir="simple_output"):
    """Simple version using only reliable dependencies"""
    os.makedirs(output_dir, exist_ok=True)
    
    print("1. Extracting text from PDF...")
    text_parts = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                text_parts.append(text.strip())
    
    full_text = "\n\n".join(text_parts)
    print(f"Extracted {len(full_text)} characters")
    
    if len(full_text) < 100:
        print("Error: Not enough text extracted")
        return
    
    print("2. Creating slides...")
    # Simple text chunking
    paragraphs = [p for p in full_text.split('\n\n') if p.strip() and len(p.strip()) > 50]
    slides = paragraphs[:8]  # First 8 paragraphs as slides
    
    print(f"Created {len(slides)} slides")
    
    # Create PowerPoint
    print("3. Creating PowerPoint...")
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    slide.shapes.title.text = "PDF to Lecture"
    slide.placeholders[1].text = "Generated from PDF"
    
    # Content slides
    for i, content in enumerate(slides, 1):
        content_slide = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide)
        slide.shapes.title.text = f"Slide {i}"
        
        # Truncate content if too long
        if len(content) > 1000:
            content = content[:1000] + "..."
        
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = content
    
    pptx_path = os.path.join(output_dir, "lecture.pptx")
    prs.save(pptx_path)
    print(f"PowerPoint saved: {pptx_path}")
    
    # Create simple images for video
    print("4. Creating slide images...")
    images = []
    for i, slide_text in enumerate(slides, 1):
        img = Image.new('RGB', (800, 600), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        # Simple text rendering
        wrapped = textwrap.fill(slide_text, width=60)
        lines = wrapped.split('\n')
        
        y = 50
        for line in lines[:15]:  # Limit to 15 lines
            draw.text((50, y), line, fill=(0, 0, 0))
            y += 30
            if y > 550:
                break
        
        img_path = os.path.join(output_dir, f"slide_{i}.png")
        img.save(img_path)
        images.append(img_path)
    
    print("5. Generating audio...")
    combined_text = " ".join(slides)
    if len(combined_text) > 4000:  # gTTS limit
        combined_text = combined_text[:4000]
    
    tts = gTTS(combined_text)
    audio_path = os.path.join(output_dir, "narration.mp3")
    tts.save(audio_path)
    print(f"Audio saved: {audio_path}")
    
    print("6. Process complete!")
    print(f"Output directory: {output_dir}")

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        simple_pdf_to_lecture(sys.argv[1])
    else:
        print("Usage: python pure_pdfplumber_app.py your_file.pdf")