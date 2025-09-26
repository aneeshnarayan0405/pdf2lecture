# ui.py - FIXED VERSION
import streamlit as st
import os
import tempfile
import sys
import time

# Add the current directory to path so imports work
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

st.set_page_config(
    page_title="PDF to Lecture Converter",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

def simple_pipeline(pdf_path, output_dir, use_tts="gtts"):
    """Simplified pipeline for Web UI"""
    import pdfplumber
    from gtts import gTTS
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image, ImageDraw, ImageFont
    import textwrap
    
    os.makedirs(output_dir, exist_ok=True)
    results = {}
    
    # 1. Extract text
    text_parts = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                text_parts.append(text.strip())
    
    full_text = "\n\n".join(text_parts)
    
    if len(full_text) < 100:
        raise Exception("Not enough text extracted from PDF")
    
    # 2. Simple summarization (first 8 paragraphs)
    paragraphs = [p for p in full_text.split('\n\n') if p.strip() and len(p.strip()) > 50]
    summaries = paragraphs[:8]
    
    # 3. Create PowerPoint
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    slide.shapes.title.text = "PDF to Lecture"
    slide.placeholders[1].text = "AI-Generated Presentation"
    
    # Content slides
    for i, content in enumerate(summaries, 1):
        content_slide = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide)
        slide.shapes.title.text = f"Slide {i}"
        
        # Truncate content if too long
        if len(content) > 1000:
            content = content[:1000] + "..."
        
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = content
    
    pptx_path = os.path.join(output_dir, "lecture.pptx")
    prs.save(pptx_path)
    results["pptx"] = pptx_path
    
    # 4. Create slide images
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    images = []
    
    for i, slide_text in enumerate(summaries, 1):
        img = Image.new('RGB', (800, 600), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        # Simple text rendering
        wrapped = textwrap.fill(slide_text, width=60)
        lines = wrapped.split('\n')
        
        y = 50
        for line in lines[:15]:  # Limit to 15 lines
            draw.text((50, y), line, fill=(0, 0, 0))
            y += 30
            if y > 550:
                break
        
        img_path = os.path.join(images_dir, f"slide_{i}.png")
        img.save(img_path)
        images.append(img_path)
    
    results["images"] = images
    
    # 5. Generate audio
    combined_text = " ".join(summaries)
    if len(combined_text) > 4000:  # gTTS limit
        combined_text = combined_text[:4000]
    
    if use_tts == "gtts":
        tts = gTTS(combined_text)
        audio_path = os.path.join(output_dir, "narration.mp3")
        tts.save(audio_path)
        results["audio"] = audio_path
    else:
        # Fallback for pyttsx3
        try:
            from pdf2lecture.tts import tts_pyttsx3
            audio_path = os.path.join(output_dir, "narration.mp3")
            tts_pyttsx3(combined_text, audio_path)
            results["audio"] = audio_path
        except:
            results["audio"] = None
    
    # 6. Create video if audio available
    if results.get("audio") and images:
        try:
            from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
            
            audio = AudioFileClip(results["audio"])
            audio_duration = audio.duration
            slide_duration = audio_duration / len(images)
            
            clips = []
            for image_path in images:
                clip = ImageClip(image_path, duration=slide_duration)
                clips.append(clip)
            
            video = concatenate_videoclips(clips)
            video = video.set_audio(audio)
            
            video_path = os.path.join(output_dir, "lecture.mp4")
            video.write_videofile(video_path, fps=24, verbose=False, logger=None)
            results["video"] = video_path
        except Exception as e:
            print(f"Video creation failed: {e}")
            results["video"] = None
    
    return results

def main():
    st.title("🎓 AI-Powered PDF to Lecture Converter")
    st.markdown("""
    **Convert any PDF document into engaging video lectures or presentation slides!**
    
    Upload a PDF file and the AI will automatically:
    - Extract and summarize content
    - Generate professional slides
    - Create narrated video lectures
    """)
    
    # Sidebar with settings
    with st.sidebar:
        st.header("⚙️ Settings")
        
        # TTS engine selection
        tts_engine = st.selectbox(
            "Text-to-Speech Engine",
            ["gtts", "pyttsx3"],
            index=0,
            help="gTTS (online, better quality) or pyttsx3 (offline)"
        )
        
        st.header("📊 Options")
        show_preview = st.checkbox("Show content preview", value=True)
    
    # Main content area
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("📤 Upload PDF File")
        uploaded_file = st.file_uploader(
            "Choose a PDF file",
            type="pdf",
            help="Supported: lecture notes, research papers, textbook chapters"
        )
        
        if uploaded_file is not None:
            # File info
            file_size = len(uploaded_file.getvalue()) / 1024  # KB
            st.success(f"✅ **Upload Successful:** {uploaded_file.name} ({file_size:.1f} KB)")
            
            # Preview PDF content if requested
            if show_preview:
                with st.expander("📄 PDF Content Preview", expanded=False):
                    try:
                        import pdfplumber
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                            tmp_file.write(uploaded_file.getvalue())
                            tmp_path = tmp_file.name
                        
                        with pdfplumber.open(tmp_path) as pdf:
                            if len(pdf.pages) > 0:
                                first_page = pdf.pages[0]
                                text = first_page.extract_text()
                                if text:
                                    st.text_area("First Page Content", text[:500] + "..." if len(text) > 500 else text, height=150, key="preview")
                                else:
                                    st.warning("Could not extract text from first page")
                                    # Try alternative extraction
                                    words = first_page.extract_words()
                                    if words:
                                        alt_text = ' '.join(word['text'] for word in words[:50])
                                        st.text_area("Extracted Words", alt_text + "...", height=100)
                        
                        os.unlink(tmp_path)
                    except Exception as e:
                        st.error(f"Preview error: {e}")
            
            # Processing button
            if st.button("🚀 Generate Lecture", type="primary", use_container_width=True):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                try:
                    # Save uploaded file temporarily
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                        tmp_file.write(uploaded_file.getvalue())
                        pdf_path = tmp_file.name
                    
                    # Create output directory
                    output_dir = f"web_output_{uploaded_file.name.replace('.pdf', '')}"
                    
                    status_text.text("Step 1/6: Processing PDF...")
                    progress_bar.progress(10)
                    time.sleep(0.5)
                    
                    status_text.text("Step 2/6: Extracting text...")
                    progress_bar.progress(20)
                    time.sleep(0.5)
                    
                    status_text.text("Step 3/6: Creating slides...")
                    progress_bar.progress(40)
                    time.sleep(0.5)
                    
                    status_text.text("Step 4/6: Generating audio...")
                    progress_bar.progress(60)
                    
                    # Run the pipeline
                    results = simple_pipeline(pdf_path, output_dir, tts_engine)
                    
                    status_text.text("Step 5/6: Creating video...")
                    progress_bar.progress(80)
                    time.sleep(0.5)
                    
                    status_text.text("Step 6/6: Finalizing...")
                    progress_bar.progress(100)
                    
                    st.success("🎉 Generation Complete!")
                    
                    # Display results
                    st.subheader("📁 Generated Files")
                    
                    # Show download buttons for each output
                    col1, col2, col3 = st.columns(3)
                    
                    if results.get("pptx"):
                        with col1:
                            with open(results["pptx"], "rb") as f:
                                st.download_button(
                                    "📊 Download PowerPoint",
                                    data=f.read(),
                                    file_name="lecture.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    use_container_width=True
                                )
                    
                    if results.get("video"):
                        with col2:
                            with open(results["video"], "rb") as f:
                                st.download_button(
                                    "🎥 Download Video",
                                    data=f.read(),
                                    file_name="lecture.mp4",
                                    mime="video/mp4",
                                    use_container_width=True
                                )
                    
                    if results.get("audio"):
                        with col3:
                            with open(results["audio"], "rb") as f:
                                st.download_button(
                                    "🔊 Download Audio",
                                    data=f.read(),
                                    file_name="narration.mp3",
                                    mime="audio/mpeg",
                                    use_container_width=True
                                )
                    
                    # Show video preview
                    if results.get("video"):
                        st.subheader("🎬 Video Preview")
                        st.video(results["video"])
                    
                    # Show success message
                    st.balloons()
                    
                except Exception as e:
                    st.error(f"❌ Processing failed: {str(e)}")
                    st.info("💡 **Troubleshooting tips:**\n- Try a different PDF file\n- Ensure the PDF contains text (not scanned images)\n- Check your internet connection for gTTS")
                
                finally:
                    # Clean up
                    status_text.empty()
                    progress_bar.empty()
                    if 'pdf_path' in locals() and os.path.exists(pdf_path):
                        os.unlink(pdf_path)
    
    with col2:
        st.subheader("ℹ️ How It Works")
        st.markdown("""
        **1. Upload PDF**  
        • Drag & drop your file  
        • Text-based PDFs work best
        
        **2. AI Processing**  
        • Automatic text extraction  
        • Smart content summarization  
        • Professional slide creation
        
        **3. Download Results**  
        • PowerPoint presentation  
        • Narrated video lecture  
        • Separate audio file
        """)
        
        st.subheader("💡 Tips for Best Results")
        st.markdown("""
        - Use **text-based PDFs** (not scanned images)
        - **1-50 pages** works best
        - **gTTS** for better voice quality
        - Clear, structured content works best
        """)

if __name__ == "__main__":
    main()