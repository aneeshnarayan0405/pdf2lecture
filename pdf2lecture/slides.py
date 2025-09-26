from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import os
import textwrap
from typing import List

def create_pptx(slide_texts: List[str], output_path: str = "lecture.pptx", 
                title: str = "Generated Lecture") -> str:
    """
    Create PowerPoint presentation from slide texts.
    
    Args:
        slide_texts: List of text content for each slide
        output_path: Output PPTX file path
        title: Presentation title
        
    Returns:
        Path to created PPTX file
    """
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Automatically generated from PDF"
    
    # Content slides
    for i, content in enumerate(slide_texts, 1):
        # Use title and content layout
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Slide title
        slide.shapes.title.text = f"Slide {i}"
        
        # Content
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = ""
        text_frame.word_wrap = True
        
        # Split content into bullet points
        paragraphs = [p.strip() for p in content.split('. ') if p.strip()]
        
        for j, paragraph in enumerate(paragraphs[:6]):  # Max 6 bullets
            if j == 0:
                p = text_frame.paragraphs[0]
                p.text = paragraph
                if not paragraph.endswith('.'):
                    p.text += '.'
            else:
                p = text_frame.add_paragraph()
                p.text = paragraph
                if not paragraph.endswith('.'):
                    p.text += '.'
                p.level = 1
            
            p.font.size = Pt(18)
            p.font.name = "Calibri"
    
    prs.save(output_path)
    return output_path

def create_slide_images(slide_texts: List[str], output_dir: str = "slide_images",
                       size: tuple = (1280, 720)) -> List[str]:
    """
    Create PNG images for each slide (for video generation).
    
    Args:
        slide_texts: List of slide contents
        output_dir: Directory to save images
        size: Image dimensions (width, height)
        
    Returns:
        List of paths to created images
    """
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []
    
    # Try to load a nice font
    try:
        # Try different possible font paths
        font_paths = [
            "/System/Library/Fonts/Helvetica.ttc",  # macOS
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",  # Linux
            "C:/Windows/Fonts/arial.ttf",  # Windows
        ]
        font = None
        for path in font_paths:
            if os.path.exists(path):
                font = ImageFont.truetype(path, 28)
                break
        if font is None:
            font = ImageFont.load_default()
    except:
        font = ImageFont.load_default()
    
    for i, text in enumerate(slide_texts, 1):
        # Create image
        img = Image.new('RGB', size, color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        # Draw slide number
        title = f"Slide {i}"
        title_bbox = draw.textbbox((0, 0), title, font=font)
        title_width = title_bbox[2] - title_bbox[0]
        draw.text(((size[0] - title_width) // 2, 30), title, font=font, fill=(0, 0, 0))
        
        # Draw content (wrapped text)
        margin = 50
        y_position = 100
        max_width = size[0] - 2 * margin
        
        # Wrap text
        wrapped_text = textwrap.fill(text, width=50)
        lines = wrapped_text.split('\n')
        
        for line in lines:
            line_bbox = draw.textbbox((0, 0), line, font=font)
            line_height = line_bbox[3] - line_bbox[1] + 5
            draw.text((margin, y_position), line, font=font, fill=(0, 0, 0))
            y_position += line_height
            
            # Stop if running out of space
            if y_position > size[1] - 50:
                draw.text((margin, y_position), "...", font=font, fill=(0, 0, 0))
                break
        
        # Save image
        filename = os.path.join(output_dir, f"slide_{i:02d}.png")
        img.save(filename)
        image_paths.append(filename)
    
    return image_paths

def add_speaker_notes(pptx_path: str, notes: List[str]) -> str:
    """
    Add speaker notes to PowerPoint presentation.
    
    Args:
        pptx_path: Path to PPTX file
        notes: List of notes (one per slide, excluding title slide)
        
    Returns:
        Path to modified PPTX
    """
    prs = Presentation(pptx_path)
    
    # Skip title slide (index 0), add notes to content slides
    for i, slide in enumerate(prs.slides[1:], 0):
        if i < len(notes):
            text_frame = slide.notes_slide.notes_text_frame
            text_frame.text = notes[i]
    
    prs.save(pptx_path)
    return pptx_path